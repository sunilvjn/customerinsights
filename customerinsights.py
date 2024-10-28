from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
import json
import os
from io import BytesIO
import glob
import logging

class TranscriptReader:
    def __init__(self, transcript_directory):
        self.transcript_directory = transcript_directory
        self.conversation_summaries = []
        
    def read_transcripts(self):
        """Read and process JSON transcript files"""
        transcripts = []
        json_files = glob.glob(os.path.join(self.transcript_directory, "*.json"))
        
        print(f"Found {len(json_files)} JSON files")
        
        for file_path in json_files[:1000]:  # Limit to 1000 files
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    conversations = json.load(file)
                    
                    for conversation in conversations:
                        if self.is_valid_conversation(conversation):
                            # Extract conversation details
                            transcript_info = {
                                'conversation_id': conversation['conversation_id'],
                                'customer_id': conversation['customer_id'],
                                'agent_name': conversation['agent']['name'],
                                'issue_type': conversation['metadata']['issue_type'],
                                'resolution_status': conversation['metadata']['resolution_status'],
                                'customer_satisfaction': conversation['metrics']['customer_satisfaction'],
                                'duration_seconds': conversation['metadata']['duration_seconds'],
                                'transcript': self.format_conversation(conversation['messages']),
                                'sentiment_summary': self.summarize_sentiment(conversation['messages']),
                                'tags': conversation['tags']
                            }
                            
                            self.conversation_summaries.append({
                                'conversation_id': conversation['conversation_id'],
                                'issue_type': conversation['metadata']['issue_type'],
                                'satisfaction': conversation['metrics']['customer_satisfaction'],
                                'resolution_status': conversation['metadata']['resolution_status'],
                                'duration': conversation['metadata']['duration_seconds'],
                                'sentiment': self.summarize_sentiment(conversation['messages'])
                            })
                            
                            transcripts.append(transcript_info)
                            
            except Exception as e:
                logging.error(f"Error reading file {file_path}: {str(e)}")
        
        print(f"Successfully processed {len(transcripts)} conversations")
        return transcripts
    
    def is_valid_conversation(self, conversation):
        """Validate if conversation has all required fields and correct data types"""
        required_fields = {
            'conversation_id': str,
            'customer_id': str,
            'agent': dict,
            'metadata': dict,
            'metrics': dict,
            'messages': list
        }
        
        try:
            return all(
                isinstance(conversation.get(field), type_)
                for field, type_ in required_fields.items()
            )
        except (AttributeError, TypeError):
            return False
    
    def format_conversation(self, messages):
        """Format messages into a readable transcript"""
        formatted_messages = []
        for message in messages:
            timestamp = message['timestamp']
            speaker = message['speaker']
            content = message['message']
            sentiment = message.get('sentiment', 'neutral')
            
            formatted_message = f"[{timestamp}] {speaker}: {content} (Sentiment: {sentiment})"
            formatted_messages.append(formatted_message)
            
        return "\n".join(formatted_messages)
    
    def summarize_sentiment(self, messages):
        """Analyze sentiment distribution in the conversation"""
        sentiment_counts = {
            'positive': 0,
            'negative': 0,
            'neutral': 0,
            'frustrated': 0,
            'worried': 0,
            'concerned': 0
        }
        
        for message in messages:
            sentiment = message.get('sentiment', 'neutral').lower()
            if sentiment in sentiment_counts:
                sentiment_counts[sentiment] += 1
        
        total_messages = len(messages)
        sentiment_percentages = {
            k: (v / total_messages) * 100 
            for k, v in sentiment_counts.items() 
            if v > 0  # Only include sentiments that appeared
        }
        
        return sentiment_percentages
    
    def get_analytics_data(self):
        """Get aggregated analytics from processed conversations"""
        if not self.conversation_summaries:
            return None
            
        analytics = {
            'issue_types': {},
            'avg_satisfaction': 0,
            'resolution_rate': 0,
            'avg_duration': 0,
            'sentiment_distribution': {}
        }
        
        total_conversations = len(self.conversation_summaries)
        
        # Calculate metrics
        for summary in self.conversation_summaries:
            # Count issue types
            issue_type = summary['issue_type']
            analytics['issue_types'][issue_type] = analytics['issue_types'].get(issue_type, 0) + 1
            
            # Add to averages
            analytics['avg_satisfaction'] += summary['satisfaction']
            analytics['avg_duration'] += summary['duration']
            
            # Count resolutions
            if summary['resolution_status'] == 'resolved':
                analytics['resolution_rate'] += 1
        
        # Finalize calculations
        analytics['avg_satisfaction'] /= total_conversations
        analytics['avg_duration'] /= total_conversations
        analytics['resolution_rate'] = (analytics['resolution_rate'] / total_conversations) * 100
        
        # Convert issue type counts to percentages
        analytics['issue_types'] = {
            k: (v / total_conversations) * 100 
            for k, v in analytics['issue_types'].items()
        }
        
        return analytics

class PresentationFormatter:
    def __init__(self):
        self.title_font_size = Pt(32)
        self.heading_font_size = Pt(24)
        self.body_font_size = Pt(14)
        self.brand_color = RGBColor(0, 112, 192)  # Blue

    def apply_slide_template(self, slide):
        """Apply consistent formatting to slide"""
        if slide.shapes.title:
            title = slide.shapes.title
            title.text_frame.paragraphs[0].font.size = self.title_font_size
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = self.brand_color

    def format_bullet_points(self, text_frame):
        """Apply consistent formatting to bullet points"""
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = self.body_font_size
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
            paragraph.level = 0

class DataVisualizer:
    def create_issue_trend_chart(self, data):
        """Create trend chart for top issues"""
        # Convert dictionary to DataFrame with proper structure
        df = pd.DataFrame(list(data.items()), columns=['Issue', 'Percentage'])
        
        plt.figure(figsize=(10, 6))
        sns.barplot(x='Issue', y='Percentage', data=df)
        plt.xticks(rotation=45)
        plt.title('Top Customer Issues by Frequency')
        plt.tight_layout()

        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        plt.close()
        return img_stream

    def create_sentiment_donut(self, sentiment_data):
        """Create donut chart for sentiment distribution"""
        plt.figure(figsize=(8, 8))
        plt.pie(sentiment_data.values(), labels=sentiment_data.keys(), autopct='%1.1f%%',
                pctdistance=0.85, hole=.5)
        plt.title('Customer Sentiment Distribution')

        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        plt.close()
        return img_stream

def analyze_with_openai(client, transcripts, analysis_type):
    """Enhanced analysis using OpenAI with structured data"""
    prompt_templates = {
        'executive_summary': """Analyze these customer service interactions and provide an executive summary.
            Include:
            1) Overall sentiment trends
            2) Key performance metrics (avg satisfaction, resolution rate)
            3) Most common issue types
            4) Notable patterns in customer behavior
            5) Agent performance insights
            Format as clear, concise bullet points.""",

        'product_improvement': """Based on these customer interactions, identify product improvement opportunities.
            Consider:
            1) Frequently reported issues
            2) Customer pain points
            3) Specific feature requests
            4) Competitive mentions
            5) Impact on customer satisfaction
            Include relevant metrics and examples.""",

        'complaints': """Analyze patterns in customer complaints from these interactions.
            Focus on:
            1) Most common complaint categories
            2) Severity distribution
            3) Resolution patterns
            4) Customer impact
            5) Satisfaction correlation
            Provide specific examples and metrics.""",

        'marketing': """Extract marketing insights from these customer interactions.
            Include:
            1) Customer segment patterns
            2) Product feature preferences
            3) Satisfaction drivers
            4) Cross-sell opportunities
            5) Customer communication preferences
            Support with specific examples and metrics.""",

        'top_issues': """Analyze the top 5 customer issues from these interactions.
            For each issue provide:
            1) Issue description and frequency
            2) Average resolution time
            3) Customer satisfaction impact
            4) Common resolution paths
            5) Improvement recommendations
            Include specific metrics and examples.""",

        'sentiment_analysis': """Provide detailed sentiment analysis of these customer interactions.
            Include:
            1) Overall sentiment distribution
            2) Sentiment trends by issue type
            3) Impact on resolution time
            4) Correlation with satisfaction
            5) Agent response effectiveness
            Support with specific metrics and examples."""
    }

    # Prepare structured data for analysis
    analytics_data = []
    for transcript in transcripts[:50]:  # Analyze 50 conversations at a time
        analytics_data.append({
            'conversation_id': transcript['conversation_id'],
            'issue_type': transcript['issue_type'],
            'satisfaction': transcript['customer_satisfaction'],
            'resolution_status': transcript['resolution_status'],
            'duration': transcript['duration_seconds'],
            'transcript': transcript['transcript']
        })

    prompt = f"{prompt_templates[analysis_type]}\n\nAnalytics Data:\n{json.dumps(analytics_data, indent=2)}"

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a customer insight analyst specialized in analyzing customer service interactions. Provide data-driven insights with specific metrics and examples."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.3
    )

    return response.choices[0].message.content

class PresentationBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.formatter = PresentationFormatter()
        self.visualizer = DataVisualizer()

    def add_title_slide(self):
        """Add main title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Customer Interaction Analysis"
        subtitle.text = "Insights & Recommendations"

        self.formatter.apply_slide_template(slide)

    def add_executive_summary(self, summary_content):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Executive Summary"
        content.text = summary_content

        self.formatter.apply_slide_template(slide)
        self.formatter.format_bullet_points(content.text_frame)

    def add_analysis_slide_with_chart(self, title, content, chart_stream, chart_position):
        """Add analysis slide with visualization"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title_shape = slide.shapes.title
        title_shape.text = title

        content_box = slide.shapes.placeholders[1]
        content_box.text = content

        chart_stream.seek(0)
        slide.shapes.add_picture(chart_stream,
                               chart_position['left'],
                               chart_position['top'],
                               chart_position['width'],
                               chart_position['height'])

        self.formatter.apply_slide_template(slide)
        self.formatter.format_bullet_points(content_box.text_frame)

def main():
    # Configuration
    TRANSCRIPT_DIRECTORY = "/transcripts"  # Replace with your transcript directory
    OPENAI_API_KEY = "your-api-key"
    
    # Check if directory exists
    if not os.path.exists(TRANSCRIPT_DIRECTORY):
        print(f"Error: Directory {TRANSCRIPT_DIRECTORY} does not exist")
        return
        
    # Initialize clients
    client = OpenAI(api_key=OPENAI_API_KEY)
    transcript_reader = TranscriptReader(TRANSCRIPT_DIRECTORY)
    
    # Read transcripts
    print("Reading transcripts...")
    transcripts = transcript_reader.read_transcripts()
    print(f"Found {len(transcripts)} transcripts")
    
    # If no transcripts found, use sample data for demonstration
    if len(transcripts) == 0:
        print("No transcripts found. Using sample data for demonstration.")
        transcripts = [
            {
                'conversation_id': 'sample1',
                'customer_id': 'cust1',
                'issue_type': 'Technical',
                'resolution_status': 'resolved',
                'customer_satisfaction': 4,
                'duration_seconds': 300,
                'transcript': 'Sample transcript'
            }
        ]

    # Initialize presentation builder
    presentation = PresentationBuilder()
    
    # Add title slide
    presentation.add_title_slide()

    # Get executive summary
    print("Generating executive summary...")
    exec_summary = analyze_with_openai(client, transcripts, "executive_summary")
    presentation.add_executive_summary(exec_summary)

    # Analyze and create slides for each category
    analyses = {
        "Product Improvement Opportunities": "product_improvement",
        "Customer Complaints Analysis": "complaints",
        "Marketing Insights": "marketing",
        "Top 5 Customer Issues": "top_issues",
        "Customer Sentiment Analysis": "sentiment_analysis"
    }

    for title, analysis_type in analyses.items():
        print(f"Analyzing {title}...")
        content = analyze_with_openai(client, transcripts, analysis_type)

        # Create appropriate visualization
        if analysis_type == "top_issues":
            issues_data = {
                "Technical Issues": 35,
                "Billing Problems": 25,
                "Product Questions": 20,
                "Account Access": 15,
                "Service Requests": 5
            }
            chart = presentation.visualizer.create_issue_trend_chart(issues_data)
        elif analysis_type == "sentiment_analysis":
            sentiment_data = {"Positive": 45, "Neutral": 30, "Negative": 25}
            chart = presentation.visualizer.create_sentiment_donut(sentiment_data)
        else:
            continue

        chart_position = {
            'left': Inches(5),
            'top': Inches(2),
            'width': Inches(4),
            'height': Inches(3)
        }
        presentation.add_analysis_slide_with_chart(title, content, chart, chart_position)

    # Save presentation
    output_file = "customer_interaction_analysis.pptx"
    presentation.prs.save(output_file)
    print(f"Presentation saved as {output
